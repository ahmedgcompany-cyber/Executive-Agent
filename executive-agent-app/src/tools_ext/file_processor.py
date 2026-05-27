"""File processor — reads and extracts content from any file format."""

import os
from pathlib import Path
from typing import Any, Optional


class FileProcessor:
    """Extract readable content from uploaded files of any format."""

    # Formats we can handle
    SUPPORTED = {
        "text":   [".txt", ".md", ".csv", ".log", ".py", ".js", ".ts", ".html",
                   ".css", ".json", ".xml", ".yaml", ".yml", ".ini", ".cfg",
                   ".sh", ".bat", ".sql", ".r", ".java", ".c", ".cpp", ".h", ".go", ".rs"],
        "pdf":    [".pdf"],
        "word":   [".docx", ".doc"],
        "excel":  [".xlsx", ".xls", ".csv"],
        "image":  [".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff"],
        "pptx":   [".pptx", ".ppt"],
        "email":  [".eml", ".msg"],
        "audio":  [".mp3", ".wav", ".m4a"],   # note-only, no transcription
    }

    @classmethod
    def can_process(cls, file_path: str) -> bool:
        ext = Path(file_path).suffix.lower()
        return any(ext in exts for exts in cls.SUPPORTED.values())

    @classmethod
    def get_category(cls, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        for cat, exts in cls.SUPPORTED.items():
            if ext in exts:
                return cat
        return "unknown"

    def process(self, file_path: str) -> dict[str, Any]:
        """Extract text content and metadata from a file.

        Args:
            file_path: Absolute or relative path to the file

        Returns:
            dict with: success, content (str), metadata (dict), category, file_name
        """
        path = Path(file_path)
        if not path.exists():
            return {"success": False, "error": f"File not found: {file_path}"}

        category = self.get_category(file_path)
        file_name = path.name
        file_size = path.stat().st_size

        try:
            if category == "text":
                return self._read_text(path, file_name, file_size)
            elif category == "pdf":
                return self._read_pdf(path, file_name, file_size)
            elif category == "word":
                return self._read_word(path, file_name, file_size)
            elif category == "excel":
                return self._read_excel(path, file_name, file_size)
            elif category == "image":
                return self._read_image(path, file_name, file_size)
            elif category == "pptx":
                return self._read_pptx(path, file_name, file_size)
            elif category == "email":
                return self._read_email(path, file_name, file_size)
            else:
                return {
                    "success": False,
                    "error": f"Unsupported file type: {path.suffix}",
                    "file_name": file_name,
                    "category": "unknown",
                }
        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to read {file_name}: {e}",
                "file_name": file_name,
                "category": category,
            }

    # ── Readers ───────────────────────────────────────────────────────

    def _read_text(self, path: Path, file_name: str, size: int) -> dict:
        """Read plain text / code files."""
        try:
            content = path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            content = path.read_text(encoding="latin-1", errors="replace")

        lines = content.count("\n")
        return {
            "success": True,
            "file_name": file_name,
            "category": "text",
            "content": content,
            "metadata": {
                "size_bytes": size,
                "lines": lines,
                "encoding": "utf-8",
            },
            "summary": f"{file_name} ({lines} lines)",
        }

    def _read_pdf(self, path: Path, file_name: str, size: int) -> dict:
        """Extract text from PDF — tries PyMuPDF, then pypdf, then pdfminer."""
        content = ""
        pages = 0
        method = ""

        # Try PyMuPDF (fitz) — best quality
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(path))
            pages = len(doc)
            parts = []
            for page in doc:
                parts.append(page.get_text())
            content = "\n".join(parts)
            method = "PyMuPDF"
        except ImportError:
            pass
        except Exception:
            pass

        # Try pypdf
        if not content:
            try:
                from pypdf import PdfReader
                reader = PdfReader(str(path))
                pages = len(reader.pages)
                parts = []
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        parts.append(t)
                content = "\n".join(parts)
                method = "pypdf"
            except ImportError:
                pass
            except Exception:
                pass

        # Try pdfminer
        if not content:
            try:
                from pdfminer.high_level import extract_text
                content = extract_text(str(path))
                method = "pdfminer"
            except ImportError:
                pass
            except Exception:
                pass

        if not content:
            return {
                "success": False,
                "error": "Could not extract PDF text. Install: pip install pymupdf",
                "file_name": file_name,
                "category": "pdf",
                "install_hint": "pip install pymupdf",
            }

        return {
            "success": True,
            "file_name": file_name,
            "category": "pdf",
            "content": content.strip(),
            "metadata": {"size_bytes": size, "pages": pages, "extractor": method},
            "summary": f"{file_name} ({pages} pages, extracted with {method})",
        }

    def _read_word(self, path: Path, file_name: str, size: int) -> dict:
        """Extract text from .docx Word documents."""
        try:
            from docx import Document
            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # Also extract table content
            tables = []
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    rows.append(" | ".join(cell.text.strip() for cell in row.cells))
                tables.append("\n".join(rows))

            content = "\n\n".join(paragraphs)
            if tables:
                content += "\n\n=== Tables ===\n" + "\n\n".join(tables)

            return {
                "success": True,
                "file_name": file_name,
                "category": "word",
                "content": content,
                "metadata": {
                    "size_bytes": size,
                    "paragraphs": len(paragraphs),
                    "tables": len(doc.tables),
                },
                "summary": f"{file_name} ({len(paragraphs)} paragraphs)",
            }
        except ImportError:
            return {
                "success": False,
                "error": "python-docx not installed. Run: pip install python-docx",
                "file_name": file_name,
                "category": "word",
                "install_hint": "pip install python-docx",
            }

    def _read_excel(self, path: Path, file_name: str, size: int) -> dict:
        """Extract data from Excel or CSV files."""
        ext = path.suffix.lower()

        if ext == ".csv":
            import csv, io
            try:
                text = path.read_text(encoding="utf-8", errors="replace")
            except Exception:
                text = path.read_text(encoding="latin-1", errors="replace")

            reader = csv.reader(io.StringIO(text))
            rows = list(reader)
            content = "\n".join(" | ".join(row) for row in rows[:500])  # max 500 rows
            return {
                "success": True,
                "file_name": file_name,
                "category": "excel",
                "content": content,
                "metadata": {"rows": len(rows), "format": "csv"},
                "summary": f"{file_name} ({len(rows)} rows, CSV)",
            }

        # xlsx / xls
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"=== Sheet: {sheet_name} ===")
                for row in ws.iter_rows(max_row=500, values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c for c in cells):
                        parts.append(" | ".join(cells))
            content = "\n".join(parts)
            return {
                "success": True,
                "file_name": file_name,
                "category": "excel",
                "content": content,
                "metadata": {"sheets": len(wb.sheetnames), "format": "xlsx"},
                "summary": f"{file_name} ({len(wb.sheetnames)} sheet(s))",
            }
        except ImportError:
            return {
                "success": False,
                "error": "openpyxl not installed. Run: pip install openpyxl",
                "file_name": file_name,
                "category": "excel",
                "install_hint": "pip install openpyxl",
            }

    def _read_image(self, path: Path, file_name: str, size: int) -> dict:
        """Return image metadata and optionally extract text via OCR."""
        meta: dict = {"size_bytes": size, "format": path.suffix.upper().lstrip(".")}

        # Try Pillow for dimensions
        try:
            from PIL import Image as PilImage
            with PilImage.open(str(path)) as img:
                meta["width"] = img.width
                meta["height"] = img.height
                meta["mode"] = img.mode
        except ImportError:
            pass
        except Exception:
            pass

        # Try OCR via pytesseract
        ocr_text = ""
        try:
            import pytesseract
            from PIL import Image as PilImage
            ocr_text = pytesseract.image_to_string(PilImage.open(str(path)))
        except ImportError:
            pass
        except Exception:
            pass

        if ocr_text.strip():
            content = f"[Image: {file_name}]\nOCR Text:\n{ocr_text.strip()}"
        else:
            dims = f"{meta.get('width', '?')}x{meta.get('height', '?')}" if "width" in meta else ""
            content = (
                f"[Image file: {file_name}"
                + (f"  {dims}" if dims else "")
                + "]\n"
                + "(No OCR text extracted — install pytesseract for text extraction)"
            )

        return {
            "success": True,
            "file_name": file_name,
            "category": "image",
            "content": content,
            "file_path": str(path.resolve()),
            "metadata": meta,
            "summary": f"{file_name} ({meta.get('width','?')}x{meta.get('height','?')} image)",
        }

    def _read_pptx(self, path: Path, file_name: str, size: int) -> dict:
        """Extract text from PowerPoint slides."""
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                slides_text.append(f"=== Slide {i} ===\n" + "\n".join(parts))

            content = "\n\n".join(slides_text)
            return {
                "success": True,
                "file_name": file_name,
                "category": "pptx",
                "content": content,
                "metadata": {"slides": len(prs.slides), "size_bytes": size},
                "summary": f"{file_name} ({len(prs.slides)} slides)",
            }
        except ImportError:
            return {
                "success": False,
                "error": "python-pptx not installed. Run: pip install python-pptx",
                "file_name": file_name,
                "category": "pptx",
                "install_hint": "pip install python-pptx",
            }

    def _read_email(self, path: Path, file_name: str, size: int) -> dict:
        """Read .eml email files."""
        try:
            import email as _email
            raw = path.read_bytes()
            msg = _email.message_from_bytes(raw)
            subject = msg.get("Subject", "")
            from_addr = msg.get("From", "")
            to_addr = msg.get("To", "")
            date = msg.get("Date", "")

            body_parts = []
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        charset = part.get_content_charset() or "utf-8"
                        try:
                            body_parts.append(part.get_payload(decode=True).decode(charset, errors="replace"))
                        except Exception:
                            pass
            else:
                charset = msg.get_content_charset() or "utf-8"
                body_parts.append(msg.get_payload(decode=True).decode(charset, errors="replace"))

            content = (
                f"Subject: {subject}\nFrom: {from_addr}\nTo: {to_addr}\nDate: {date}\n\n"
                + "\n".join(body_parts)
            )
            return {
                "success": True,
                "file_name": file_name,
                "category": "email",
                "content": content,
                "metadata": {"subject": subject, "from": from_addr, "date": date},
                "summary": f"Email: {subject} (from {from_addr})",
            }
        except Exception as e:
            return {"success": False, "error": str(e), "file_name": file_name, "category": "email"}

    # ── Static helper for building context block ──────────────────────

    @staticmethod
    def build_context_block(file_result: dict, max_chars: int = 8000) -> str:
        """Build a readable context block to prepend to a user message.

        Args:
            file_result: Output from process()
            max_chars: Max characters of content to include

        Returns:
            Formatted string block describing the file contents
        """
        if not file_result.get("success"):
            return f"[File: {file_result.get('file_name', '?')} — could not be read: {file_result.get('error', '')}]\n"

        file_name = file_result.get("file_name", "file")
        category = file_result.get("category", "")
        summary = file_result.get("summary", "")
        content = file_result.get("content", "")

        truncated = len(content) > max_chars
        snippet = content[:max_chars] + ("\n... [truncated]" if truncated else "")

        lines = [
            f"=== Attached File: {file_name} ===",
            f"Type: {category.upper()}  |  {summary}",
            "--- Content ---",
            snippet,
            "=== End of File ===",
        ]
        return "\n".join(lines)
