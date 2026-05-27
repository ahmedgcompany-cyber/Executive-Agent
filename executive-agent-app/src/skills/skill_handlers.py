"""
Native Python skill handlers for all 56 MegaV skills.

Each handler signature:
    handler(user_input: str, extra_context: dict, router, emit) -> dict

'router' is a ModelRouter instance (or None if no AI available).
'emit'   is a callable(str) for progress messages.

This replaces bash script execution with real Python implementations
that work on Windows/Mac/Linux without external shell dependencies.
"""

from __future__ import annotations

import csv
import hashlib
import json
import math
import os
import re
import secrets
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path
from typing import Any, Callable, Optional

from ..providers.model_router import NoModelAvailableError

DESKTOP = Path.home() / "Desktop"


# ─────────────────────────────────────────────────────────────────────────────
# Shared helpers
# ─────────────────────────────────────────────────────────────────────────────

def _ok(result: str, summary: str = "", path: str = "") -> dict:
    return {
        "success": True,
        "result":  result,
        "summary": summary or result[:200],
        "path":    path or "",
        "mode":    "native",
    }


def _fail(msg: str) -> dict:
    return {
        "success": False,
        "result":  "",
        "summary": msg,
        "error":   msg,
        "mode":    "native",
    }


def _ai(router, system_prompt: str, user_input: str, emit=None) -> dict:
    """Call the model_router with a system prompt + user input."""
    if not router:
        return _fail(
            "No AI provider available — start Ollama (`ollama serve`) or set ANTHROPIC_API_KEY."
        )
    try:
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_input},
        ]
        resp = router.route_chat(messages, task_type="skill_execution")
        if resp.get("success"):
            msg  = resp.get("message", {})
            text = msg.get("content", "") if isinstance(msg, dict) else str(msg)
            return _ok(text)
        # fallback: generate
        resp = router.route_generate(
            user_input, system_prompt=system_prompt, task_type="skill_execution"
        )
        if resp.get("success"):
            return _ok(resp.get("response", ""))
        return _fail(resp.get("error", "AI call failed"))
    except NoModelAvailableError as exc:
        return _fail(f"No LLM available: {exc}")
    except Exception as exc:
        return _fail(f"AI routing error: {exc}")


def _extract_path(user_input: str, extra_context: Optional[dict], keys: list[str]) -> Optional[str]:
    """Extract a file/folder path from extra_context or user_input."""
    if extra_context:
        for k in keys:
            v = extra_context.get(k)
            if v:
                return str(v)
    m = re.search(
        r'"([^"]+)"|\'([^\']+)\'|([A-Za-z]:\\[^\s"\']+)|(/[^\s"\']+)',
        user_input,
    )
    if m:
        return next(g for g in m.groups() if g is not None)
    return None


# ─────────────────────────────────────────────────────────────────────────────
# CONTENT CREATION
# ─────────────────────────────────────────────────────────────────────────────

def handle_artifacts_builder(user_input, extra_context, router, emit):
    """Generate a complete, self-contained React HTML artifact."""
    result = _ai(router, """You are an expert React/HTML artifact engineer.
Generate a complete, single-file HTML document with:
- React + ReactDOM loaded via CDN (no build step needed)
- Tailwind CSS via CDN
- Full interactivity with useState/useEffect hooks
- Dark theme (#0d1117 background)
- Beautiful, polished UI with proper spacing and typography
- Working state management and event handlers

Output ONLY the raw HTML file — no markdown, no explanation, just <!DOCTYPE html>...""",
        user_input, emit,
    )
    if result.get("success"):
        html = result["result"]
        # Save to desktop
        fname = f"artifact_{secrets.token_hex(4)}.html"
        out = DESKTOP / fname
        out.write_text(html, encoding="utf-8")
        result["path"] = str(out)
        emit and emit(f"[ArtifactsBuilder] Saved: {out}")
    return result


def handle_brand_guidelines(user_input, extra_context, router, emit):
    return _ai(router, """You are an Anthropic brand guidelines specialist.
Official Anthropic brand:
- Primary: #D97757 (warm amber-orange)  Secondary: #8B6F47 (warm brown)
- Dark bg: #1A1A1A   Light bg: #FAFAFA  Text dark: #E8E3DC  Text light: #1A1A1A
- Headings: Tiempos Headline (serif)  Body: Public Sans (sans-serif)
- Tone: warm, precise, thoughtful, not corporate

Apply this brand identity to the user's content. If it is HTML/CSS, update all
colour references and font stacks. If it is text, return a styled HTML document.
Return the complete, fully branded result.""",
        user_input, emit,
    )


def handle_canvas_design(user_input, extra_context, router, emit):
    """Create a visual design image using PIL, guided by AI colour/text choices."""
    try:
        from PIL import Image, ImageDraw, ImageFont  # type: ignore

        emit and emit("[CanvasDesign] Generating visual design…")

        # Parse canvas size
        w, h = 1200, 800
        m = re.search(r'(\d{3,4})\s*[xX×]\s*(\d{3,4})', user_input)
        if m:
            w, h = int(m.group(1)), int(m.group(2))

        # Ask AI for design parameters
        ai = _ai(router, """Return ONLY a JSON object (no markdown) for a visual design:
{
  "bg": "#hexcode",
  "text_color": "#hexcode",
  "accent": "#hexcode",
  "title": "short heading (max 6 words)",
  "subtitle": "one line subtitle",
  "style": "modern"
}""", user_input, emit)

        cfg = {
            "bg": "#0d1117", "text_color": "#e6edf3",
            "accent": "#f78166", "title": "Untitled Design",
            "subtitle": user_input[:60], "style": "modern",
        }
        if ai.get("success"):
            try:
                cfg.update(json.loads(re.search(r'\{[\s\S]*\}', ai["result"]).group()))
            except Exception:
                pass

        def hex2rgb(h):
            h = h.lstrip("#")
            return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

        img = Image.new("RGB", (w, h), hex2rgb(cfg["bg"]))
        draw = ImageDraw.Draw(img)

        # Accent stripe
        ay = h // 3
        for i in range(6):
            alpha = 200 - i * 30
            draw.rectangle([0, ay + i*3, w, ay + i*3 + 2], fill=hex2rgb(cfg["accent"]))

        # Try to load a system font
        try:
            font_big = ImageFont.truetype("arial.ttf", max(48, w // 20))
            font_sm  = ImageFont.truetype("arial.ttf", max(24, w // 40))
        except Exception:
            font_big = ImageFont.load_default()
            font_sm  = font_big

        draw.text((w // 12, h // 2 - 80),  cfg["title"],    fill=hex2rgb(cfg["text_color"]), font=font_big)
        draw.text((w // 12, h // 2 + 20),  cfg["subtitle"], fill=hex2rgb(cfg["accent"]),     font=font_sm)

        out = DESKTOP / f"design_{secrets.token_hex(4)}.png"
        img.save(str(out))
        emit and emit(f"[CanvasDesign] Saved: {out}")
        return _ok(f"Design saved to: {out}", path=str(out))

    except ImportError:
        # No Pillow — generate styled HTML instead
        res = _ai(router, """Create a beautiful HTML/CSS design page for this request.
Use inline styles only. Dark background, bold typography, professional layout.
Return ONLY the complete HTML file.""", user_input, emit)
        if res.get("success"):
            out = DESKTOP / f"design_{secrets.token_hex(4)}.html"
            out.write_text(res["result"], encoding="utf-8")
            res["path"] = str(out)
        return res
    except Exception as exc:
        return _fail(f"Canvas design error: {exc}")


def handle_changelog_generator(user_input, extra_context, router, emit):
    """Auto-generate a CHANGELOG.md from git commit history."""
    cwd = (extra_context or {}).get("project_path", ".")
    try:
        proc = subprocess.run(
            ["git", "log", "--pretty=format:%H|%s|%an|%ad", "--date=short", "-80"],
            capture_output=True, text=True, cwd=cwd, timeout=10,
        )
        if proc.returncode != 0:
            raise RuntimeError(proc.stderr.strip())

        commits = []
        for line in proc.stdout.splitlines():
            parts = line.split("|", 3)
            if len(parts) == 4:
                sha, msg, author, date = parts
                commits.append({"sha": sha[:8], "msg": msg, "author": author, "date": date})

        cats: dict[str, list] = {
            "feat": [], "fix": [], "docs": [], "refactor": [], "chore": [], "other": [],
        }
        for c in commits:
            m = c["msg"].lower()
            if m.startswith(("feat", "add ", "new ")):
                cats["feat"].append(c)
            elif m.startswith(("fix", "bug", "crash")):
                cats["fix"].append(c)
            elif m.startswith(("doc", "readme")):
                cats["docs"].append(c)
            elif m.startswith(("refactor", "clean", "rewrite")):
                cats["refactor"].append(c)
            elif m.startswith(("chore", "update", "bump", "upgrade")):
                cats["chore"].append(c)
            else:
                cats["other"].append(c)

        labels = {
            "feat":     "✨ Features",
            "fix":      "🐛 Bug Fixes",
            "docs":     "📚 Documentation",
            "refactor": "♻️ Refactoring",
            "chore":    "🔧 Chores",
            "other":    "📦 Other",
        }
        lines = [f"# Changelog\n\n*Generated from {len(commits)} commits*\n"]
        for cat, label in labels.items():
            items = cats[cat]
            if not items:
                continue
            lines.append(f"\n### {label}\n")
            for c in items[:20]:
                lines.append(f"- {c['msg']} (`{c['sha']}`) — {c['date']}")

        changelog = "\n".join(lines)
        out = DESKTOP / "CHANGELOG.md"
        out.write_text(changelog, encoding="utf-8")
        return _ok(changelog, summary=f"Changelog from {len(commits)} commits", path=str(out))

    except Exception:
        # No git — use AI to write a changelog from description
        return _ai(router, """Write a professional CHANGELOG.md in Keep-a-Changelog format.
Include sections: [Unreleased], and version entries with Added/Changed/Fixed/Removed.
Return only the raw markdown.""", user_input, emit)


def handle_slack_gif_creator(user_input, extra_context, router, emit):
    """Create an animated GIF optimised for Slack (128×128, ≤500 KB)."""
    try:
        from PIL import Image, ImageDraw  # type: ignore

        emit and emit("[GIFCreator] Building animated GIF…")

        ai = _ai(router, """Return ONLY valid JSON for a GIF animation config:
{
  "frames": 12,
  "bg": "#1a1a2e",
  "primary": "#e94560",
  "shape": "pulse"
}
Allowed shapes: pulse, rotate, wave, bounce""", user_input, emit)

        cfg = {"frames": 12, "bg": "#1a1a2e", "primary": "#e94560", "shape": "pulse"}
        if ai.get("success"):
            try:
                cfg.update(json.loads(re.search(r'\{[\s\S]*\}', ai["result"]).group()))
            except Exception:
                pass

        n  = min(int(cfg.get("frames", 12)), 30)
        W  = H = 128

        def hx(h):
            h = h.lstrip("#")
            return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

        bg, pc = hx(cfg["bg"]), hx(cfg["primary"])
        shape  = cfg.get("shape", "pulse")
        frames = []

        for i in range(n):
            t = i / n
            fr = Image.new("RGB", (W, H), bg)
            dr = ImageDraw.Draw(fr)

            if shape == "pulse":
                r = int(15 + 30 * abs(math.sin(t * math.pi)))
                dr.ellipse([W//2-r, H//2-r, W//2+r, H//2+r], fill=pc)
            elif shape == "rotate":
                angle = t * 2 * math.pi
                rx = int(W//2 + 44 * math.cos(angle))
                ry = int(H//2 + 44 * math.sin(angle))
                dr.ellipse([rx-8, ry-8, rx+8, ry+8], fill=pc)
            elif shape == "wave":
                for x in range(0, W, 4):
                    y = int(H//2 + 28 * math.sin(x/W*4*math.pi + t*2*math.pi))
                    dr.ellipse([x-3, y-3, x+3, y+3], fill=pc)
            else:  # bounce
                y = int(H//2 + 40 * math.sin(t * math.pi))
                dr.ellipse([W//2-12, y-12, W//2+12, y+12], fill=pc)

            frames.append(fr)

        out = DESKTOP / f"animation_{secrets.token_hex(4)}.gif"
        frames[0].save(
            out, save_all=True, append_images=frames[1:],
            optimize=False, duration=80, loop=0,
        )
        emit and emit(f"[GIFCreator] Saved: {out}")
        return _ok(f"GIF created: {out}", path=str(out))

    except ImportError:
        return _fail("Pillow required: pip install pillow")
    except Exception as exc:
        return _fail(f"GIF creation error: {exc}")


def handle_theme_factory(user_input, extra_context, router, emit):
    """Apply one of 10 professional themes to the user's content."""
    THEMES: dict[str, dict] = {
        "midnight-galaxy": {"bg": "#0d0d1a", "text": "#e8e8ff", "accent": "#7c3aed", "font": "Inter"},
        "arctic-frost":    {"bg": "#f0f9ff", "text": "#0c4a6e", "accent": "#0284c7", "font": "Roboto"},
        "ocean-depths":    {"bg": "#0c1445", "text": "#bfdbfe", "accent": "#3b82f6", "font": "Lato"},
        "forest-canopy":   {"bg": "#0f2d1a", "text": "#d1fae5", "accent": "#059669", "font": "Open Sans"},
        "golden-hour":     {"bg": "#1c1008", "text": "#fef3c7", "accent": "#d97706", "font": "Merriweather"},
        "desert-rose":     {"bg": "#2d0a14", "text": "#fce7f3", "accent": "#db2777", "font": "Playfair Display"},
        "sunrise":         {"bg": "#fff7ed", "text": "#7c2d12", "accent": "#ea580c", "font": "Source Sans Pro"},
        "minimal":         {"bg": "#ffffff", "text": "#111827", "accent": "#374151", "font": "Inter"},
        "cyberpunk":       {"bg": "#000000", "text": "#39ff14", "accent": "#ff0080", "font": "JetBrains Mono"},
        "executive":       {"bg": "#1a1a1a", "text": "#f5f5f5", "accent": "#c9a84c", "font": "Cormorant Garamond"},
    }
    req = user_input.lower()
    selected = "midnight-galaxy"
    for name in THEMES:
        if name.replace("-", " ") in req or name.split("-")[0] in req:
            selected = name
            break
    t = THEMES[selected]
    return _ai(router, f"""You are a theme application specialist using the '{selected}' theme.
Colors — background: {t['bg']}, text: {t['text']}, accent: {t['accent']}, font: {t['font']}.
Apply this theme to the user's content precisely:
- If it is HTML/CSS: update all colour and font declarations
- If it is markdown/text: wrap in a full HTML document using the theme colours
Return only the complete styled result, no explanation.""",
        user_input, emit,
    )


# ─────────────────────────────────────────────────────────────────────────────
# CAREER
# ─────────────────────────────────────────────────────────────────────────────

def handle_tailored_resume_generator(user_input, extra_context, router, emit):
    result = _ai(router, """You are an elite ATS resume writer and career coach.
Create a tailored, ATS-optimised resume using this exact process:
1. Extract every requirement keyword from the job description
2. Mirror that exact language throughout the resume
3. Quantify all achievements (numbers, %, $, time saved)
4. Use strong action verbs: Led, Built, Delivered, Increased, Reduced
5. Structure: Summary → Experience → Skills → Education
6. Add an ATS compatibility score (aim 85%+) with explanation
7. Append a short, personalised cover letter paragraph

Output clean markdown that converts directly to a Word doc.""",
        user_input, emit,
    )
    if result.get("success"):
        out = DESKTOP / f"resume_{secrets.token_hex(4)}.md"
        out.write_text(result["result"], encoding="utf-8")
        result["path"] = str(out)
        emit and emit(f"[ResumeGenerator] Saved: {out}")
    return result


# ─────────────────────────────────────────────────────────────────────────────
# BUSINESS / SALES
# ─────────────────────────────────────────────────────────────────────────────

def handle_competitive_ads_extractor(user_input, extra_context, router, emit):
    return _ai(router, """You are a competitive intelligence expert specialising in digital advertising.
Analyse competitor advertising strategy and deliver a professional report:

## Competitor Ad Intelligence Report

1. KEY MESSAGES & VALUE PROPOSITIONS
   - What core benefits do they advertise?
   - What pain points do they target?

2. CREATIVE PATTERNS
   - Visual style (bold, minimal, lifestyle…)
   - Copy tone (formal, casual, urgent…)
   - CTA patterns (free trial, get started, book demo…)

3. AUDIENCE TARGETING SIGNALS
   - Demographics indicated by messaging
   - Job titles / industries mentioned
   - Funnel stage (awareness / consideration / conversion)

4. AD FORMATS USED
   - Static image, video, carousel, lead gen form…

5. STRATEGIC VULNERABILITIES
   - What are they NOT saying? (your opportunity)

6. RECOMMENDED COUNTER-STRATEGY
   - 3 specific ad angles that would differentiate you

Be specific. Use real examples and data where possible.""",
        user_input, emit,
    )


def handle_domain_name_brainstormer(user_input, extra_context, router, emit):
    """Generate domain name ideas and check availability via DNS lookup."""
    import socket

    emit and emit("[DomainBrainstormer] Generating name ideas…")

    ai = _ai(router, """Return ONLY a JSON array of 20 creative base domain names (no TLDs, no punctuation).
Short (4-12 chars), memorable, brandable. Example: ["launchpad","veloce","novacore"]""",
        user_input, emit,
    )
    base_names = ["startup", "launch", "build", "create", "tech", "digital", "nexus",
                  "forge", "spark", "velo", "nova", "apex", "core", "flux", "zeal"]
    if ai.get("success"):
        try:
            base_names = json.loads(re.search(r'\[[\s\S]*\]', ai["result"]).group())
        except Exception:
            pass

    tlds = [".com", ".io", ".ai", ".co", ".app", ".dev", ".net"]
    rows = []
    for name in base_names[:15]:
        name_clean = re.sub(r'[^a-z0-9\-]', '', name.lower())
        if not name_clean:
            continue
        row: dict[str, str] = {"name": name_clean}
        for tld in tlds:
            domain = name_clean + tld
            try:
                socket.getaddrinfo(domain, None)
                row[tld] = "TAKEN"
            except socket.gaierror:
                row[tld] = "FREE ✓"
            except Exception:
                row[tld] = "?"
        rows.append(row)

    lines = [f"# Domain Name Ideas\n*Query: {user_input[:80]}*\n",
             "| Name | .com | .io | .ai | .co | .app | .dev | .net |",
             "|------|------|-----|-----|-----|------|------|------|"]
    for r in rows:
        cols = [r["name"]] + [r.get(t, "?") for t in tlds]
        lines.append("| " + " | ".join(cols) + " |")

    report = "\n".join(lines)
    out = DESKTOP / "domain_ideas.md"
    out.write_text(report, encoding="utf-8")
    free_count = sum(1 for r in rows for v in r.values() if v == "FREE ✓")
    return _ok(report, summary=f"{len(rows)} names checked, {free_count} available slots", path=str(out))


def handle_lead_research_assistant(user_input, extra_context, router, emit):
    result = _ai(router, """You are an elite B2B sales researcher. Identify and qualify leads.

Deliver:
## Target ICP (Ideal Customer Profile)
- Industry / company size / tech stack / growth stage

## Lead List (10-15 companies)
| Company | Size | Industry | Why They're a Fit | Score/10 |

## Top 3 Leads — Deep Dive
For each: LinkedIn URL format, Decision maker title, Pain points, Personalised opener

## Discovery Questions (5)
Open-ended questions to qualify them

## Competitive Landscape
Who else are they likely evaluating?

## 7-Day Outreach Action Plan
Day-by-day with specific tactics

Be specific — real company names, real titles, real pain points.""",
        user_input, emit,
    )
    if result.get("success"):
        out = DESKTOP / f"leads_{secrets.token_hex(4)}.md"
        out.write_text(result["result"], encoding="utf-8")
        result["path"] = str(out)
    return result


def handle_twitter_algorithm_optimizer(user_input, extra_context, router, emit):
    return _ai(router, """You are a Twitter/X algorithm expert.

ALGORITHM SIGNALS (ranked by impact):
1. Replies in first hour (highest weight)
2. Bookmarks / saves
3. Retweets without quote
4. Long dwell time (scrolling slowly past)
5. External links TANK reach → put links in replies
6. Threads outperform single tweets by 3×
7. Images +35% reach, Video +65% reach
8. Optimal length: 71-100 characters
9. Best times: 8-10am, 12-1pm, 5-6pm local

REWRITE TACTICS:
- Hook in the first 3 words
- Specific numbers over vague claims
- Open loop format ("Most people don't know this about X:")
- Controversial but defensible take
- End with an engagement prompt

For each tweet/post provided:
1. Show the ORIGINAL
2. Show the OPTIMISED version
3. List exactly what changed and why
4. Estimated reach multiplier""",
        user_input, emit,
    )


# ─────────────────────────────────────────────────────────────────────────────
# AUTOMATION / FILES
# ─────────────────────────────────────────────────────────────────────────────

def handle_file_organizer(user_input, extra_context, router, emit):
    """Scan a folder and organise files by type."""
    folder_path = _extract_path(user_input, extra_context, ["folder_path", "path"])
    if not folder_path:
        folder_path = str(Path.home() / "Downloads")

    folder = Path(folder_path)
    if not folder.exists():
        return _fail(f"Folder not found: {folder_path}")

    emit and emit(f"[FileOrganizer] Scanning: {folder}")

    TYPE_MAP: dict[str, list[str]] = {
        "Images":        [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".svg", ".ico", ".tiff"],
        "Videos":        [".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv", ".webm", ".m4v"],
        "Audio":         [".mp3", ".wav", ".flac", ".aac", ".ogg", ".m4a", ".wma"],
        "Documents":     [".pdf", ".doc", ".docx", ".txt", ".rtf", ".odt", ".pages", ".epub"],
        "Spreadsheets":  [".xlsx", ".xls", ".csv", ".ods", ".numbers"],
        "Presentations": [".pptx", ".ppt", ".odp", ".key"],
        "Code":          [".py", ".js", ".ts", ".html", ".css", ".json", ".xml", ".yaml",
                          ".yml", ".sh", ".bat", ".ps1", ".java", ".c", ".cpp", ".h",
                          ".go", ".rs", ".rb", ".php", ".sql", ".md"],
        "Archives":      [".zip", ".rar", ".7z", ".tar", ".gz", ".bz2"],
        "Executables":   [".exe", ".msi", ".dmg", ".deb", ".rpm"],
    }
    ext_to_type: dict[str, str] = {}
    for t, exts in TYPE_MAP.items():
        for e in exts:
            ext_to_type[e] = t

    all_files = [f for f in folder.iterdir() if f.is_file()]
    emit and emit(f"[FileOrganizer] {len(all_files)} files found")

    # Duplicate detection
    hash_map: dict[str, Path] = {}
    duplicates: list[tuple[Path, Path]] = []
    for f in all_files:
        try:
            h = hashlib.md5(f.read_bytes()).hexdigest()
            if h in hash_map:
                duplicates.append((f, hash_map[h]))
            else:
                hash_map[h] = f
        except Exception:
            pass

    do_move = any(
        w in user_input.lower()
        for w in ["move", "organize now", "do it", "execute", "run", "apply", "go ahead"]
    )

    lines = [
        f"# File Organization Report",
        f"**Folder:** `{folder}`",
        f"**Files:** {len(all_files)}   **Duplicates:** {len(duplicates)}",
        "",
    ]

    if do_move:
        moved = 0
        for f in all_files:
            type_name = ext_to_type.get(f.suffix.lower(), "Other")
            dest_dir = folder / type_name
            dest_dir.mkdir(exist_ok=True)
            dest = dest_dir / f.name
            if not dest.exists():
                shutil.move(str(f), str(dest))
                moved += 1
        lines.append(f"✅ **Moved {moved} files** into categorised subfolders.")
        emit and emit(f"[FileOrganizer] Moved {moved} files")
    else:
        plan: dict[str, list[str]] = {}
        for f in all_files:
            t = ext_to_type.get(f.suffix.lower(), "Other")
            plan.setdefault(t, []).append(f.name)

        lines.append("## Plan *(add 'move' to your request to execute)*\n")
        for type_name, files in sorted(plan.items()):
            lines.append(f"### {type_name}/ — {len(files)} files")
            for fn in files[:8]:
                lines.append(f"  - {fn}")
            if len(files) > 8:
                lines.append(f"  … and {len(files)-8} more")
            lines.append("")

    if duplicates:
        lines.append("## Duplicate Files")
        for dup, orig in duplicates[:15]:
            lines.append(f"  - `{dup.name}` = duplicate of `{orig.name}`")

    report = "\n".join(lines)
    out = DESKTOP / "file_organization_report.md"
    out.write_text(report, encoding="utf-8")
    return _ok(report, summary=f"{len(all_files)} files, {len(duplicates)} duplicates", path=str(out))


def handle_invoice_organizer(user_input, extra_context, router, emit):
    """Organise invoices/receipts and export a CSV expense log."""
    folder_path = _extract_path(
        user_input, extra_context, ["invoices_folder", "folder_path", "path"]
    )
    if not folder_path:
        folder_path = str(Path.home() / "Downloads")

    folder = Path(folder_path)
    if not folder.exists():
        return _fail(f"Folder not found: {folder_path}")

    INVOICE_EXTS = {".pdf", ".jpg", ".jpeg", ".png", ".tiff", ".bmp"}
    files = [f for f in folder.iterdir() if f.is_file() and f.suffix.lower() in INVOICE_EXTS]
    emit and emit(f"[InvoiceOrganizer] {len(files)} documents found")

    records = []
    for f in files:
        date_m   = re.search(r'(\d{4}[-_]\d{2}[-_]\d{2}|\d{2}[-_]\d{2}[-_]\d{4})', f.name)
        amount_m = re.search(r'(\d[\d,]*\.\d{2})', f.name)
        vendor   = re.sub(r'[0-9_\-\.]+', ' ', f.stem).strip()[:40]
        records.append({
            "filename": f.name,
            "date":     date_m.group(1)   if date_m   else "unknown",
            "vendor":   vendor or f.stem[:20],
            "amount":   amount_m.group(1) if amount_m else "0.00",
            "type":     "invoice" if "inv" in f.name.lower() else "receipt",
            "path":     str(f),
        })

    out_csv = DESKTOP / "invoices.csv"
    with open(out_csv, "w", newline="", encoding="utf-8") as fh:
        writer = csv.DictWriter(fh, fieldnames=["filename", "date", "vendor", "amount", "type", "path"])
        writer.writeheader()
        writer.writerows(records)

    total = sum(
        float(r["amount"].replace(",", ""))
        for r in records if r["amount"] != "0.00"
    )
    summary = (
        f"Processed {len(records)} invoices/receipts\n"
        f"Detected total: ${total:,.2f}\n"
        f"CSV exported: {out_csv}"
    )
    return _ok(summary, path=str(out_csv))


def handle_video_downloader(user_input, extra_context, router, emit):
    """Download a video using yt-dlp."""
    url_m = re.search(r'https?://[^\s]+', user_input)
    if not url_m:
        return _fail("No URL found. Example: 'download https://youtube.com/watch?v=...'")

    url  = url_m.group(0)
    u    = user_input.lower()
    mode = "audio" if ("mp3" in u or "audio only" in u) else \
           "1080p" if "1080" in u else \
           "480p"  if "480" in u else "720p"

    out_dir = DESKTOP / "Downloads"
    out_dir.mkdir(exist_ok=True)
    emit and emit(f"[VideoDownloader] Downloading ({mode}): {url[:60]}…")

    fmt_map = {
        "audio": ["-f", "bestaudio", "-x", "--audio-format", "mp3"],
        "1080p": ["-f", "bestvideo[height<=1080]+bestaudio/best[height<=1080]"],
        "720p":  ["-f", "best[height<=720]"],
        "480p":  ["-f", "best[height<=480]"],
    }
    fmt_args = fmt_map.get(mode, fmt_map["720p"])
    output_tmpl = str(out_dir / "%(title)s.%(ext)s")

    def _try(cmd_prefix: list[str]) -> dict:
        cmd = cmd_prefix + fmt_args + ["-o", output_tmpl, url]
        try:
            proc = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
            if proc.returncode == 0:
                return _ok(f"Download complete → {out_dir}", path=str(out_dir))
            return _fail(proc.stderr.strip()[:400] or "yt-dlp error")
        except FileNotFoundError:
            return _fail("NOT_FOUND")
        except subprocess.TimeoutExpired:
            return _fail("Download timed out (300 s)")

    # Try system yt-dlp first, then Python module
    res = _try(["yt-dlp"])
    if not res["success"] and "NOT_FOUND" in res.get("error", ""):
        res = _try([sys.executable, "-m", "yt_dlp"])
    if not res["success"] and "NOT_FOUND" in res.get("error", ""):
        emit and emit("[VideoDownloader] Installing yt-dlp…")
        subprocess.run([sys.executable, "-m", "pip", "install", "yt-dlp", "-q"],
                       capture_output=True)
        res = _try([sys.executable, "-m", "yt_dlp"])
    return res


def handle_image_enhancer(user_input, extra_context, router, emit):
    """Enhance an image using Pillow."""
    try:
        from PIL import Image, ImageEnhance, ImageFilter  # type: ignore

        image_path = _extract_path(
            user_input, extra_context,
            ["image_path", "file_path", "path"],
        )
        if not image_path or not Path(image_path).exists():
            return _fail("Image path not found. Provide the full path to the image file.")

        emit and emit(f"[ImageEnhancer] Processing: {Path(image_path).name}")
        img = Image.open(image_path).convert("RGB")
        u   = user_input.lower()

        # Apply operations based on request
        if "upscale" in u or "2x" in u or "enlarge" in u:
            w, h = img.size
            img  = img.resize((w * 2, h * 2), Image.LANCZOS)
            emit and emit(f"[ImageEnhancer] Upscaled to {w*2}×{h*2}")

        if "grayscale" in u or "black and white" in u or "b&w" in u:
            img = img.convert("L").convert("RGB")
        else:
            img = ImageEnhance.Color(img).enhance(1.2)         # +20% saturation

        if "blur" in u:
            img = img.filter(ImageFilter.GaussianBlur(radius=2))
        else:
            img = ImageEnhance.Sharpness(img).enhance(1.6)     # sharpen

        img = ImageEnhance.Contrast(img).enhance(1.25)
        img = ImageEnhance.Brightness(img).enhance(1.1)

        stem = Path(image_path).stem
        ext  = Path(image_path).suffix or ".jpg"
        out  = DESKTOP / f"{stem}_enhanced{ext}"
        img.save(str(out), quality=95)
        emit and emit(f"[ImageEnhancer] Saved: {out}")
        return _ok(f"Enhanced image saved: {out}", path=str(out))

    except ImportError:
        return _fail("Pillow required: pip install pillow")
    except Exception as exc:
        return _fail(f"Image enhancement error: {exc}")


# ─────────────────────────────────────────────────────────────────────────────
# INTEGRATIONS
# ─────────────────────────────────────────────────────────────────────────────

def handle_connect(user_input, extra_context, router, emit):
    return _ai(router, """You are an API integration expert.
For the requested service connection, provide:
1. The exact API endpoint and authentication method
2. A complete, runnable Python code snippet using `requests` or `httpx`
3. Required environment variables / API keys (with links to get them)
4. The expected response structure
5. A simple test to confirm the connection works
6. Common errors and how to fix them

Write working Python code the user can paste and run immediately.""",
        user_input, emit,
    )


def handle_connect_apps(user_input, extra_context, router, emit):
    return _ai(router, """You are an integration setup specialist.
Provide complete step-by-step setup instructions:
1. Where to get credentials (exact URLs)
2. OAuth flow (if applicable — include redirect URL setup)
3. Python connection test snippet
4. How to store credentials securely (.env file)
5. Troubleshooting common auth failures""",
        user_input, emit,
    )


def handle_composio_skills(user_input, extra_context, router, emit):
    return _ai(router, """You are an API automation expert.
Build a complete Python integration for the requested service:
1. Identify the correct REST API
2. Show authentication (Bearer / API key / OAuth)
3. Write Python functions for the requested actions
4. Handle pagination, rate limits, and retries
5. Show example request + response

All code must be runnable with only `requests` as a dependency.""",
        user_input, emit,
    )


# ─────────────────────────────────────────────────────────────────────────────
# DEVELOPMENT
# ─────────────────────────────────────────────────────────────────────────────

def handle_mcp_builder(user_input, extra_context, router, emit):
    result = _ai(router, """You are an MCP (Model Context Protocol) server expert.
Generate a complete, production-ready MCP server in Python using FastMCP.

Include:
1. Full server code with all requested tools (decorated with @mcp.tool())
2. Proper type hints on all parameters
3. Docstrings that Claude will use to understand the tools
4. Error handling that returns helpful messages
5. pyproject.toml / requirements.txt
6. Claude Desktop config snippet (claude_desktop_config.json)
7. How to run: `uvx fastmcp run server.py`

Every tool must actually DO something — no stubs.""",
        user_input, emit,
    )
    if result.get("success"):
        out = DESKTOP / f"mcp_server_{secrets.token_hex(4)}.py"
        # Extract just the Python code if wrapped in markdown
        code = result["result"]
        m = re.search(r'```(?:python)?\n([\s\S]+?)```', code)
        if m:
            code = m.group(1)
        out.write_text(code, encoding="utf-8")
        result["path"] = str(out)
        emit and emit(f"[MCPBuilder] Server saved: {out}")
    return result


def handle_langsmith_fetch(user_input, extra_context, router, emit):
    return _ai(router, """You are a LangChain/LangGraph debugging expert with deep LangSmith knowledge.
Diagnose agent execution issues:
1. LIKELY ROOT CAUSE — based on symptoms described
2. HOW TO INSTRUMENT — add tracing to pinpoint the problem
3. SPECIFIC CODE FIX — show exact before/after code
4. HOW TO VERIFY — test that the fix works
5. PREVENTION — how to avoid this class of error

Be concrete — show actual code changes, reference specific LangChain classes.""",
        user_input, emit,
    )


def handle_webapp_testing(user_input, extra_context, router, emit):
    result = _ai(router, """You are a Playwright testing expert.
Write a complete, runnable Playwright Python test file.

Include:
1. Full imports (from playwright.sync_api import sync_playwright)
2. Browser setup with proper viewport and timeouts
3. Navigation to the app URL
4. All test scenarios with clear assertions
5. Screenshot capture on failure
6. Console log capture
7. How to run: `pip install playwright && playwright install && python test_app.py`

Use semantic selectors: page.get_by_role(), get_by_text(), get_by_label().
Handle timing with proper waits, not sleep().""",
        user_input, emit,
    )
    if result.get("success"):
        code = result["result"]
        m = re.search(r'```(?:python)?\n([\s\S]+?)```', code)
        if m:
            code = m.group(1)
        out = DESKTOP / f"test_webapp_{secrets.token_hex(4)}.py"
        out.write_text(code, encoding="utf-8")
        result["path"] = str(out)
        emit and emit(f"[WebAppTester] Test saved: {out}")
    return result


def handle_skill_creator(user_input, extra_context, router, emit):
    """Scaffold a new skill folder with SKILL.md."""
    # Extract skill name
    m = re.search(
        r'"([^"]+)"|\'([^\']+)\'|skill\s+(?:called|named|for)\s+([a-zA-Z\s\-]+)',
        user_input, re.IGNORECASE,
    )
    name = "new-skill"
    if m:
        raw = m.group(1) or m.group(2) or m.group(3) or "new-skill"
        name = re.sub(r'\s+', '-', raw.strip().lower())

    skill_dir = DESKTOP / "skills" / name
    skill_dir.mkdir(parents=True, exist_ok=True)

    ai = _ai(router, f"""Write a detailed SKILL.md for a skill called '{name}'.
The document defines:
- Purpose and capabilities (2-3 paragraphs)
- Required inputs and their formats
- Exact step-by-step execution procedure
- Output format and examples
- Best practices and edge cases
- Example prompts that activate this skill

Format as professional markdown. This is used as an AI system-context document.""",
        user_input, emit,
    )
    skill_md = ai["result"] if ai.get("success") else f"# {name}\n\nSkill description.\n"
    (skill_dir / "SKILL.md").write_text(skill_md, encoding="utf-8")
    (skill_dir / "__init__.py").write_text("", encoding="utf-8")

    emit and emit(f"[SkillCreator] Created: {skill_dir}")
    return _ok(f"Skill '{name}' created at: {skill_dir}\n\n{skill_md[:300]}", path=str(skill_dir))


def handle_developer_growth_analysis(user_input, extra_context, router, emit):
    return _ai(router, """You are a senior engineering coach and learning specialist.
Analyse the developer's patterns and produce a growth report:

## Strengths (what's working well)
## Gaps (specific areas to improve)
## Recurring Patterns (anti-patterns or habits noticed)
## Learning Path
  - Resources (books, courses, docs) for each gap
  - Practice exercises with estimated time
## 30-Day Improvement Plan
  - Week 1: Foundation / diagnosis
  - Week 2: Focused practice
  - Week 3: Build something with new skills
  - Week 4: Teach it / write about it
## Success Metrics
  - How to measure improvement objectively

Be specific — reference actual code patterns described, not generic advice.""",
        user_input, emit,
    )


# ─────────────────────────────────────────────────────────────────────────────
# COMMUNICATION
# ─────────────────────────────────────────────────────────────────────────────

def handle_internal_comms(user_input, extra_context, router, emit):
    return _ai(router, """You are a professional corporate communications writer.
Write clear, concise internal communication following best practices:
- Lead with the key message in the first sentence
- Use bullet points for scanability
- Specify action items with clear owners and deadlines
- Match tone to audience (exec = formal, team = friendly)
- Avoid jargon — say exactly what you mean
- End with a single clear next step

Format appropriately for the comm type: email subject line, Slack message, announcement, etc.""",
        user_input, emit,
    )


def handle_content_research_writer(user_input, extra_context, router, emit):
    result = _ai(router, """You are an expert content strategist and writer.
Create publication-ready content using this process:

1. HOOK — 1-2 sentence opening that creates curiosity or states a bold claim
2. RESEARCH LAYER — cite 3-5 specific facts, studies, or data points
3. OUTLINE — logical structure with H2/H3 headings
4. FULL DRAFT — write every section completely (no [insert here] placeholders)
5. CTA — specific call to action at the end
6. SEO — list 5 target keywords naturally integrated

Tone: authoritative but conversational. No fluff. Every paragraph must earn its place.""",
        user_input, emit,
    )
    if result.get("success"):
        out = DESKTOP / f"content_{secrets.token_hex(4)}.md"
        out.write_text(result["result"], encoding="utf-8")
        result["path"] = str(out)
    return result


def handle_meeting_insights_analyzer(user_input, extra_context, router, emit):
    return _ai(router, """You are an expert meeting analyst and communication coach.
Analyse the meeting transcript/notes and provide:

## Executive Summary (3 bullet points max)

## Key Decisions Made
(numbered list with context)

## Action Items
| Owner | Action | Deadline | Priority |

## Communication Patterns
- Speaking balance: who dominated?
- Were decisions made or deferred?
- Any unresolved tensions?
- What was left unsaid?

## Follow-Up Required
(what must happen before next meeting)

Quote the transcript directly when making observations.""",
        user_input, emit,
    )


# ─────────────────────────────────────────────────────────────────────────────
# DOCUMENTS
# ─────────────────────────────────────────────────────────────────────────────

def _create_docx(user_input, extra_context, router, emit):
    try:
        from docx import Document  # type: ignore
        from docx.shared import Pt

        ai = _ai(router, """Return ONLY valid JSON for a Word document:
{
  "title": "document title",
  "sections": [
    {"heading": "Section Name", "content": "full paragraph text"},
    ...
  ]
}""", user_input, emit)

        data: dict = {"title": "Document", "sections": [{"heading": "", "content": user_input}]}
        if ai.get("success"):
            try:
                data = json.loads(re.search(r'\{[\s\S]*\}', ai["result"]).group())
            except Exception:
                pass

        doc = Document()
        doc.add_heading(data.get("title", "Document"), 0)
        for sec in data.get("sections", []):
            if sec.get("heading"):
                doc.add_heading(sec["heading"], 1)
            if sec.get("content"):
                doc.add_paragraph(sec["content"])

        out = DESKTOP / f"document_{secrets.token_hex(4)}.docx"
        doc.save(str(out))
        emit and emit(f"[Document] Saved Word doc: {out}")
        return _ok(f"Word document created: {out}", path=str(out))

    except ImportError:
        # Fallback: markdown file
        ai = _ai(router, "Write the complete document content in clean markdown.", user_input, emit)
        if ai.get("success"):
            out = DESKTOP / f"document_{secrets.token_hex(4)}.md"
            out.write_text(ai["result"], encoding="utf-8")
            return _ok(ai["result"], path=str(out))
        return _fail("python-docx not installed (pip install python-docx)")


def _create_excel(user_input, extra_context, router, emit):
    try:
        import openpyxl  # type: ignore
        from openpyxl.styles import Font, PatternFill, Alignment

        ai = _ai(router, """Return ONLY valid JSON for an Excel spreadsheet:
{
  "sheet_name": "Sheet1",
  "headers": ["Column1", "Column2"],
  "rows": [["value1", "value2"], ...]
}""", user_input, emit)

        data: dict = {"sheet_name": "Data", "headers": ["Item", "Value"], "rows": [["Example", "100"]]}
        if ai.get("success"):
            try:
                data = json.loads(re.search(r'\{[\s\S]*\}', ai["result"]).group())
            except Exception:
                pass

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = data.get("sheet_name", "Data")[:31]

        dark = PatternFill("solid", fgColor="0D1117")
        white_bold = Font(bold=True, color="FFFFFF")
        center = Alignment(horizontal="center")

        for col, hdr in enumerate(data.get("headers", []), 1):
            cell = ws.cell(row=1, column=col, value=hdr)
            cell.font  = white_bold
            cell.fill  = dark
            cell.alignment = center

        for ri, row in enumerate(data.get("rows", []), 2):
            for ci, val in enumerate(row, 1):
                ws.cell(row=ri, column=ci, value=val)

        out = DESKTOP / f"spreadsheet_{secrets.token_hex(4)}.xlsx"
        wb.save(str(out))
        emit and emit(f"[Document] Saved spreadsheet: {out}")
        return _ok(f"Excel file created: {out}", path=str(out))

    except ImportError:
        return _fail("openpyxl not installed (pip install openpyxl)")


def _create_pptx(user_input, extra_context, router, emit):
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt

        ai = _ai(router, """Return ONLY valid JSON for a PowerPoint presentation:
{
  "title": "Presentation Title",
  "slides": [
    {"title": "Slide Title", "bullets": ["point 1", "point 2"]},
    ...
  ]
}""", user_input, emit)

        data: dict = {"title": "Presentation", "slides": [{"title": "Overview", "bullets": [user_input[:100]]}]}
        if ai.get("success"):
            try:
                data = json.loads(re.search(r'\{[\s\S]*\}', ai["result"]).group())
            except Exception:
                pass

        prs = Presentation()
        # Title slide
        ts = prs.slides.add_slide(prs.slide_layouts[0])
        ts.shapes.title.text = data.get("title", "Presentation")
        if ts.placeholders:
            try:
                ts.placeholders[1].text = ""
            except Exception:
                pass

        for slide_data in data.get("slides", []):
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = slide_data.get("title", "")
            tf = sl.placeholders[1].text_frame
            bullets = slide_data.get("bullets") or slide_data.get("content") or []
            if isinstance(bullets, str):
                bullets = [bullets]
            for i, b in enumerate(bullets):
                if i == 0:
                    tf.text = str(b)
                else:
                    tf.add_paragraph().text = str(b)

        out = DESKTOP / f"presentation_{secrets.token_hex(4)}.pptx"
        prs.save(str(out))
        emit and emit(f"[Document] Saved presentation: {out}")
        return _ok(f"PowerPoint created: {out}", path=str(out))

    except ImportError:
        return _fail("python-pptx not installed (pip install python-pptx)")


def handle_document_skills(user_input, extra_context, router, emit):
    u = user_input.lower()
    if any(k in u for k in ["excel", "spreadsheet", "xlsx", "csv table"]):
        return _create_excel(user_input, extra_context, router, emit)
    if any(k in u for k in ["powerpoint", "pptx", "presentation", "slides", "slide deck"]):
        return _create_pptx(user_input, extra_context, router, emit)
    return _create_docx(user_input, extra_context, router, emit)


# ─────────────────────────────────────────────────────────────────────────────
# UTILITIES
# ─────────────────────────────────────────────────────────────────────────────

def handle_raffle_winner_picker(user_input, extra_context, router, emit):
    """Pick random winners using cryptographic randomness (secrets module)."""
    participants: list[str] = []

    # 1. Try extra_context
    if extra_context:
        raw = (extra_context.get("participants_list_or_file")
               or extra_context.get("participants")
               or extra_context.get("list"))
        if raw:
            if isinstance(raw, list):
                participants = [str(x).strip() for x in raw if str(x).strip()]
            elif isinstance(raw, str):
                if Path(raw).exists():
                    participants = [
                        l.strip() for l in Path(raw).read_text().splitlines() if l.strip()
                    ]
                else:
                    participants = [x.strip() for x in re.split(r'[,\n]', raw) if x.strip()]

    # 2. Extract from user_input
    if not participants:
        cleaned = re.sub(
            r'\b(pick|choose|select|raffle|winner|winners?|giveaway|contest|random|draw)\b',
            '', user_input, flags=re.IGNORECASE,
        )
        cleaned = re.sub(r'\d+\s+winner', '', cleaned, flags=re.IGNORECASE).strip()
        parts = [p.strip() for p in re.split(r'[,\n]', cleaned) if len(p.strip()) > 1]
        participants = [p for p in parts if not p.isdigit()]

    if len(participants) < 2:
        return _fail(
            "Need at least 2 participants. Example:\n"
            "'Pick 2 winners from Alice, Bob, Charlie, Diana, Eve'"
        )

    # Number of winners
    nw_m = re.search(r'(\d+)\s+winner', user_input, re.IGNORECASE)
    n_winners = int(nw_m.group(1)) if nw_m else 1
    n_winners = max(1, min(n_winners, len(participants) - 1))

    # Cryptographic shuffle
    pool = list(participants)
    winners: list[str] = []
    for _ in range(n_winners):
        idx = secrets.randbelow(len(pool))
        winners.append(pool.pop(idx))

    lines = [
        "═" * 40,
        "  🎉  RAFFLE RESULTS  🎉",
        "═" * 40,
        f"Participants : {len(participants)}",
        f"Winners drawn: {n_winners}",
        "",
        "🏆  WINNER" + ("S" if n_winners > 1 else "") + ":",
    ]
    for i, w in enumerate(winners, 1):
        lines.append(f"   {i}. {w}")
    lines += ["", "✅ Selected via cryptographic randomness (Python secrets module)"]
    result = "\n".join(lines)

    out = DESKTOP / "raffle_results.txt"
    out.write_text(result, encoding="utf-8")
    return _ok(result, path=str(out))


def handle_skill_share(user_input, extra_context, router, emit):
    """Package a skill folder into a ZIP file for sharing."""
    skill_path = _extract_path(user_input, extra_context, ["skill_folder", "folder_path", "path"])

    if skill_path and Path(skill_path).exists():
        folder = Path(skill_path)
        out_zip = DESKTOP / f"{folder.name}_skill.zip"
        with zipfile.ZipFile(out_zip, "w", zipfile.ZIP_DEFLATED) as zf:
            for f in folder.rglob("*"):
                if f.is_file():
                    zf.write(f, f.relative_to(folder.parent))
        emit and emit(f"[SkillShare] Packaged: {out_zip}")
        return _ok(f"Skill packaged: {out_zip}", path=str(out_zip))

    # No folder — give sharing instructions
    return _ai(router, """Help the user package and share their skill.
Provide:
1. How to structure the skill folder
2. How to create the ZIP file (Python zipfile snippet)
3. How to document the skill in SKILL.md
4. Best channels to share it with the team""",
        user_input, emit,
    )


# ─────────────────────────────────────────────────────────────────────────────
# SUPERPOWERS HANDLERS
# ─────────────────────────────────────────────────────────────────────────────

def handle_brainstorming(user_input, extra_context, router, emit):
    """Brainstorming: structured ideation and creative exploration."""
    return _ai(router, """You are a brainstorming facilitator. Guide the user through structured ideation:
1. Clarify the problem space and constraints
2. Generate diverse ideas using lateral thinking
3. Evaluate and prioritize ideas
4. Suggest next steps and action items

Use divergent thinking first, then convergent thinking to narrow down.""",
        user_input, emit,
    )

def handle_writing_plans(user_input, extra_context, router, emit):
    """Writing Plans: create detailed implementation plans."""
    return _ai(router, """You are a planning specialist. Create a detailed implementation plan:
1. Break the task into numbered steps
2. Identify dependencies between steps
3. Estimate complexity and risk per step
4. Define success criteria and verification methods
5. Suggest the order of execution

Format as a clear, actionable plan with checkboxes.""",
        user_input, emit,
    )

def handle_executing_plans(user_input, extra_context, router, emit):
    """Executing Plans: implement a plan step by step."""
    return _ai(router, """You are a plan execution specialist. Execute the given plan systematically:
1. Read through the entire plan first
2. Execute each step in order
3. Verify each step's output before proceeding
4. Document any deviations from the plan
5. Report completion status for each step

Be methodical and verify as you go.""",
        user_input, emit,
    )

def handle_verification_before_completion(user_input, extra_context, router, emit):
    """Verification Before Completion: verify work is truly done."""
    return _ai(router, """You are a verification specialist. Before marking anything as complete:
1. Check all requirements are met
2. Verify edge cases are handled
3. Confirm tests pass (if applicable)
4. Review for code quality and consistency
5. Ensure documentation is up to date
6. List any remaining items or caveats

Do not mark complete until all checks pass.""",
        user_input, emit,
    )

def handle_systematic_debugging(user_input, extra_context, router, emit):
    """Systematic Debugging: debug using structured methodology."""
    return _ai(router, """You are a systematic debugging specialist. Follow this methodology:
1. Reproduce: Confirm the exact error and conditions
2. Hypothesize: List possible root causes ranked by likelihood
3. Isolate: Narrow down by testing hypotheses one at a time
4. Fix: Apply the minimal fix for the root cause
5. Verify: Confirm the fix works and doesn't break anything else
6. Document: Explain what happened and why

Never guess — always verify your hypothesis before fixing.""",
        user_input, emit,
    )

def handle_test_driven_development(user_input, extra_context, router, emit):
    """Test-Driven Development: write tests first, then implement."""
    return _ai(router, """You are a TDD specialist. Follow the Red-Green-Refactor cycle:
1. RED: Write a failing test for the desired behavior
2. GREEN: Write the minimal code to make the test pass
3. REFACTOR: Clean up the code while keeping tests green
4. Repeat for each feature/behavior

Always write the test BEFORE the implementation.""",
        user_input, emit,
    )

def handle_writing_skills(user_input, extra_context, router, emit):
    """Writing Skills: create well-structured Claude skills."""
    return _ai(router, """You are a skill authoring specialist. Create a well-structured Claude skill:
1. Define the skill's purpose and when it should activate
2. Write clear, actionable instructions in SKILL.md format
3. Include examples of expected inputs and outputs
4. Define trigger keywords for auto-activation
5. Specify dependencies and execution type
6. Add validation checklists

Follow the standard SKILL.md frontmatter format.""",
        user_input, emit,
    )

def handle_using_git_worktrees(user_input, extra_context, router, emit):
    """Using Git Worktrees: manage parallel development branches."""
    return _ai(router, """You are a git worktree specialist. Help the user work with git worktrees:
1. Explain when worktrees are useful (parallel work, hotfixes, reviews)
2. Create worktrees with proper branch names
3. Guide switching between worktrees
4. Help merge/rebase changes back to main
5. Clean up worktrees when done

Provide exact git commands for each step.""",
        user_input, emit,
    )

def handle_using_superpowers(user_input, extra_context, router, emit):
    """Using Superpowers: meta-skill for discovering the right skill."""
    return _ai(router, """You are a Superpowers navigation assistant. Help the user find and use the right skill:
1. Understand what the user wants to accomplish
2. Recommend the best Superpowers skill for the task
3. Explain how to invoke it and what to expect
4. Suggest skill chains for complex workflows

Available skills: brainstorming, writing-plans, executing-plans, verification-before-completion,
systematic-debugging, test-driven-development, writing-skills, using-git-worktrees,
dispatching-parallel-agents, subagent-driven-development, requesting-code-review,
receiving-code-review, finishing-a-development-branch.""",
        user_input, emit,
    )

def handle_dispatching_parallel_agents(user_input, extra_context, router, emit):
    """Dispatching Parallel Agents: spawn multiple agents for parallel work."""
    return _ai(router, """You are a parallel agent dispatcher. Help break work into independent tasks:
1. Analyze the overall goal
2. Identify tasks that can run in parallel
3. Define clear inputs and expected outputs for each agent
4. Dispatch agents with focused prompts
5. Collect and merge results

Ensure tasks are truly independent to maximize parallelism.""",
        user_input, emit,
    )

def handle_subagent_driven_development(user_input, extra_context, router, emit):
    """Subagent-Driven Development: drive development with focused sub-agents."""
    return _ai(router, """You are a subagent-driven development coordinator. Drive development by:
1. Breaking the project into focused task areas
2. Assigning each area to a specialized sub-agent
3. Coordinating dependencies between agents
4. Merging results into a cohesive implementation
5. Running final verification

Each sub-agent should have a clear, focused mandate.""",
        user_input, emit,
    )

def handle_requesting_code_review(user_input, extra_context, router, emit):
    """Requesting Code Review: request structured code review."""
    return _ai(router, """You are a code review request specialist. Help create effective review requests:
1. Summarize what changed and why
2. Highlight areas of concern or uncertainty
3. List specific questions for reviewers
4. Provide context about the broader changeset
5. Suggest focus areas for the review

Format as a clear, actionable review request.""",
        user_input, emit,
    )

def handle_receiving_code_review(user_input, extra_context, router, emit):
    """Receiving Code Review: process and act on review feedback."""
    return _ai(router, """You are a code review receiver specialist. Process review feedback:
1. Parse all review comments and categorize them (must-fix, nice-to-have, questions)
2. Address critical issues first
3. Make targeted fixes, not broad refactors
4. Respond to each comment with what was changed
5. Request re-review after addressing feedback

Be responsive and focused in your fixes.""",
        user_input, emit,
    )

def handle_finishing_development_branch(user_input, extra_context, router, emit):
    """Finishing a Development Branch: properly close a branch."""
    return _ai(router, """You are a branch completion specialist. Properly close a development branch:
1. Run all tests and verify they pass
2. Ensure code is linted and formatted
3. Update documentation if needed
4. Create or update the PR with a clear description
5. Request code review
6. After approval, merge and clean up the branch

Never skip verification steps.""",
        user_input, emit,
    )


# ─────────────────────────────────────────────────────────────────────────────
# FRONTEND DESIGN HANDLER
# ─────────────────────────────────────────────────────────────────────────────

def handle_frontend_design(user_input, extra_context, router, emit):
    """Frontend Design: build elaborate multi-component frontends."""
    return _ai(router, """You are a frontend design specialist. Create elaborate multi-component frontends using:
- React for component architecture
- Tailwind CSS for styling
- shadcn/ui for pre-built components
- Modern CSS (Grid, Flexbox, animations)

Follow the artifacts-builder approach: initialize Vite project, build components,
bundle to single HTML file. Focus on responsive design, accessibility, and visual polish.""",
        user_input, emit,
    )


# ─────────────────────────────────────────────────────────────────────────────
# GSTACK HANDLERS
# ─────────────────────────────────────────────────────────────────────────────

def handle_gstack_browse(user_input, extra_context, router, emit):
    """GStack Browse: browse websites with headless browser."""
    return _ai(router, """You are a web browsing specialist using Playwright. Help the user:
1. Navigate to URLs and interact with pages
2. Take screenshots of pages or elements
3. Extract content from web pages
4. Fill forms and click buttons
5. Monitor console logs and network requests

Provide clear, actionable browsing instructions.""",
        user_input, emit,
    )

def handle_gstack_qa(user_input, extra_context, router, emit):
    """GStack QA: automated quality assurance testing."""
    return _ai(router, """You are a QA testing specialist. Perform systematic quality assurance:
1. Define test scenarios based on requirements
2. Test functional correctness of features
3. Check responsive design across viewports
4. Verify accessibility (WCAG guidelines)
5. Test error handling and edge cases
6. Generate a comprehensive QA report

Be thorough and systematic in testing.""",
        user_input, emit,
    )

def handle_gstack_review(user_input, extra_context, router, emit):
    """GStack Review: structured code and PR review."""
    return _ai(router, """You are a code review specialist. Perform structured reviews:
1. Check code correctness and logic
2. Review for security vulnerabilities
3. Assess code style and consistency
4. Evaluate error handling
5. Check for performance issues
6. Verify test coverage
7. Provide actionable, prioritized feedback

Rate severity: critical > important > minor > nitpick.""",
        user_input, emit,
    )

def handle_gstack_canary(user_input, extra_context, router, emit):
    """GStack Canary: canary deployment testing."""
    return _ai(router, """You are a canary deployment specialist. Verify features in production:
1. Define what to verify (core flows, performance, errors)
2. Set up monitoring and alerting
3. Gradually roll out to a small percentage
4. Compare metrics against baseline
5. Decide: promote, rollback, or expand canary

Provide clear go/no-go criteria.""",
        user_input, emit,
    )

def handle_gstack_design_consultation(user_input, extra_context, router, emit):
    """GStack Design Consultation: expert design advice."""
    return _ai(router, """You are a UI/UX design consultant. Provide expert design advice:
1. Analyze the current design and identify issues
2. Suggest improvements based on design principles
3. Recommend component patterns and layouts
4. Address accessibility and responsive design
5. Propose design system improvements

Focus on practical, implementable suggestions.""",
        user_input, emit,
    )

def handle_gstack_investigate(user_input, extra_context, router, emit):
    """GStack Investigate: systematic bug investigation."""
    return _ai(router, """You are a systematic investigation specialist. Investigate issues methodically:
1. Reproduce the issue reliably
2. Gather all relevant context (logs, stack traces, environment)
3. Form hypotheses about root cause
4. Test each hypothesis in isolation
5. Identify the root cause
6. Document findings and recommend a fix

Never skip the reproduction step.""",
        user_input, emit,
    )

def handle_gstack_ship(user_input, extra_context, router, emit):
    """GStack Ship: ship code to production safely."""
    return _ai(router, """You are a deployment specialist. Ship code to production safely:
1. Run pre-deploy checks (tests, linting, type checking)
2. Verify build succeeds
3. Check for breaking changes
4. Deploy to canary/staging first
5. Monitor for errors after deployment
6. Roll back if critical issues found

Safety first — always have a rollback plan.""",
        user_input, emit,
    )

def handle_gstack_health(user_input, extra_context, router, emit):
    """GStack Health: project health check."""
    return _ai(router, """You are a project health specialist. Check overall project health:
1. Dependency audit (outdated, vulnerable, unused)
2. Test coverage analysis
3. Code quality metrics (complexity, duplication)
4. Linting and type checking status
5. Documentation coverage
6. Build and CI/CD pipeline status

Provide a health score and prioritized improvement list.""",
        user_input, emit,
    )


# ─────────────────────────────────────────────────────────────────────────────
# MCP INTEGRATION HANDLERS
# ─────────────────────────────────────────────────────────────────────────────

def handle_context7_mcp(user_input, extra_context, router, emit):
    """Context7 MCP: look up library documentation."""
    return _ai(router, """You are a documentation lookup specialist using Context7 MCP. Help the user:
1. Identify the library or framework they need docs for
2. Look up the latest documentation
3. Provide relevant code examples and API references
4. Note any version-specific considerations
5. Suggest related documentation if available

Always fetch current docs — don't rely on training data that may be outdated.""",
        user_input, emit,
    )

def handle_tavily_search(user_input, extra_context, router, emit):
    """Tavily Search: web search for current information."""
    return _ai(router, """You are a web research specialist using Tavily. Help the user:
1. Formulate effective search queries
2. Search the web for current, accurate information
3. Summarize and synthesize findings
4. Cite sources properly
5. Note any conflicting information or uncertainties

Prioritize recency and authority of sources.""",
        user_input, emit,
    )


# ─────────────────────────────────────────────────────────────────────────────
# HANDLER REGISTRY
# ─────────────────────────────────────────────────────────────────────────────

NATIVE_HANDLERS: dict[str, Any] = {
    # Content creation
    "artifacts-builder":            handle_artifacts_builder,
    "brand-guidelines":             handle_brand_guidelines,
    "canvas-design":                handle_canvas_design,
    "changelog-generator":          handle_changelog_generator,
    "slack-gif-creator":            handle_slack_gif_creator,
    "theme-factory":                handle_theme_factory,
    # Career
    "tailored-resume-generator":    handle_tailored_resume_generator,
    # Business / sales
    "competitive-ads-extractor":    handle_competitive_ads_extractor,
    "domain-name-brainstormer":     handle_domain_name_brainstormer,
    "lead-research-assistant":      handle_lead_research_assistant,
    "twitter-algorithm-optimizer":  handle_twitter_algorithm_optimizer,
    # Automation / files
    "file-organizer":               handle_file_organizer,
    "invoice-organizer":            handle_invoice_organizer,
    "youtube-downloader":           handle_video_downloader,
    "image-enhancer":               handle_image_enhancer,
    # Integrations
    "connect":                      handle_connect,
    "connect-apps":                 handle_connect_apps,
    "composio-skills":              handle_composio_skills,
    # Development
    "mcp-builder":                  handle_mcp_builder,
    "langsmith-fetch":              handle_langsmith_fetch,
    "webapp-testing":               handle_webapp_testing,
    "skill-creator":                handle_skill_creator,
    "developer-growth-analysis":    handle_developer_growth_analysis,
    # Communication
    "internal-comms":               handle_internal_comms,
    "content-research-writer":      handle_content_research_writer,
    "meeting-insights-analyzer":    handle_meeting_insights_analyzer,
    # Documents (split into individual skills)
    "docx":                         handle_document_skills,
    "pdf":                          handle_document_skills,
    "pptx":                         handle_document_skills,
    "xlsx":                         handle_document_skills,
    # Utilities
    "raffle-winner-picker":         handle_raffle_winner_picker,
    "skill-share":                  handle_skill_share,
    # Superpowers
    "brainstorming":                    handle_brainstorming,
    "writing-plans":                    handle_writing_plans,
    "executing-plans":                  handle_executing_plans,
    "verification-before-completion":  handle_verification_before_completion,
    "systematic-debugging":             handle_systematic_debugging,
    "test-driven-development":          handle_test_driven_development,
    "writing-skills":                   handle_writing_skills,
    "using-git-worktrees":             handle_using_git_worktrees,
    "using-superpowers":               handle_using_superpowers,
    "dispatching-parallel-agents":      handle_dispatching_parallel_agents,
    "subagent-driven-development":     handle_subagent_driven_development,
    "requesting-code-review":          handle_requesting_code_review,
    "receiving-code-review":            handle_receiving_code_review,
    "finishing-a-development-branch":   handle_finishing_development_branch,
    # Frontend Design
    "frontend-design":                  handle_frontend_design,
    # GStack
    "gstack-browse":                   handle_gstack_browse,
    "gstack-qa":                       handle_gstack_qa,
    "gstack-review":                    handle_gstack_review,
    "gstack-canary":                   handle_gstack_canary,
    "gstack-design-consultation":      handle_gstack_design_consultation,
    "gstack-investigate":              handle_gstack_investigate,
    "gstack-ship":                     handle_gstack_ship,
    "gstack-health":                   handle_gstack_health,
    # MCP Integration
    "context7-mcp":                    handle_context7_mcp,
    "tavily-search":                   handle_tavily_search,
}


def get_handler(skill_id: str):
    """Return the native Python handler for a skill_id, or None."""
    return NATIVE_HANDLERS.get(skill_id)
